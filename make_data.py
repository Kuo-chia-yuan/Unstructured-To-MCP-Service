import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

# 確保資料夾存在
os.makedirs("data/raw", exist_ok=True)

# ==========================================
# 1. 產生豐富內容的瑞鼎 IC 規格說明書 (PDF)
# ==========================================
class AdvancedPDF(FPDF):
    def header(self):
        self.set_font("Arial", 'B', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, "Raydium Semiconductor Corp. - Strictly Confidential", 0, 1, 'R')
        self.ln(5)

    def footer(self):
        self.set_y(-15)
        self.set_font("Arial", 'I', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f"Page {self.page_no()}", 0, 0, 'C')

pdf = AdvancedPDF()
pdf.set_auto_page_break(auto=True, margin=15)

# --- 第一頁：封面與規格概述 ---
pdf.add_page()
pdf.set_font("Arial", 'B', 22)
pdf.set_text_color(0, 51, 102) # 科技深藍色
pdf.cell(0, 15, "Raydium Next-Gen Display Driver IC", 0, 1, 'L')
pdf.set_font("Arial", 'B', 14)
pdf.set_text_color(102, 102, 102)
pdf.cell(0, 10, "Technical Specification Sheet - v2.4.0 (2026)", 0, 1, 'L')
pdf.ln(10)

pdf.set_font("Arial", 'B', 12)
pdf.set_text_color(0, 0, 0)
pdf.cell(0, 10, "1. System Overview", 0, 1, 'L')
pdf.set_font("Arial", size=10)
pdf.multi_cell(0, 6, 
    "The RD-AI-2026 is a premium High-Definition Display Driver IC (DDIC) designed with an "
    "embedded Edge-AI inference engine. It supports real-time super-resolution, AI-driven motion "
    "blur reduction, and dynamic power optimization for OLED and high-refresh-rate mobile screens. "
    "By utilizing Model Context Protocol (MCP) server bindings, the hardware configuration can be "
    "dynamically mapped and simulated through high-level software agents."
)
pdf.ln(5)

# 插入模擬暫存器表格 (測試 Unstructured 提取表格的能力)
pdf.set_font("Arial", 'B', 11)
pdf.cell(0, 10, "Table 1.1: Electrical & Operating Parameters", 0, 1, 'L')
pdf.set_font("Arial", 'B', 10)
# 表格標頭
pdf.cell(60, 7, "Parameter", 1, 0, 'C')
pdf.cell(40, 7, "Min Value", 1, 0, 'C')
pdf.cell(40, 7, "Max Value", 1, 0, 'C')
pdf.cell(50, 7, "Interface / Protocol", 1, 1, 'C')
# 表格內容
pdf.set_font("Arial", size=10)
parameters = [
    ("Core Operating Voltage (VDD)", "0.75V", "1.25V", "Internal DC-DC"),
    ("I/O Supply Voltage (IOVDD)", "1.65V", "1.95V", "MIPI DSI-2 / SPI"),
    ("Maximum Clock Frequency", "400 MHz", "1.2 GHz", "System PLL"),
    ("Standby Power Consumption", "12 uA", "45 uA", "Deep Sleep Mode")
]
for param, p_min, p_max, p_int in parameters:
    pdf.cell(60, 6, param, 1, 0, 'L')
    pdf.cell(40, 6, p_min, 1, 0, 'C')
    pdf.cell(40, 6, p_max, 1, 0, 'C')
    pdf.cell(50, 6, p_int, 1, 1, 'L')

# --- 第二頁：暫存器配置與偵錯日誌 ---
pdf.add_page()
pdf.set_font("Arial", 'B', 12)
pdf.cell(0, 10, "2. Register Map & Configuration Tooling", 0, 1, 'L')
pdf.set_font("Arial", size=10)
pdf.multi_cell(0, 6, 
    "Host communication is implemented via standard 4-wire SPI or MIPI Command Mode. "
    "The built-in AI accelerator is initialized by flashing the firmware into the adapter memory zone. "
    "Below is the recommended initialization sequence matrix for the AI upscaling pipe:"
)
pdf.ln(4)

# 暫存器數據
pdf.set_font("Arial", 'B', 10)
pdf.cell(30, 7, "Address", 1, 0, 'C')
pdf.cell(40, 7, "Register Name", 1, 0, 'C')
pdf.cell(30, 7, "Default Hex", 1, 0, 'C')
pdf.cell(90, 7, "Functional Description", 1, 1, 'C')

pdf.set_font("Arial", size=9)
registers = [
    ("0x00A0", "AI_EN_REG", "0x00000000", "Bit[0]: Enable inference engine layer calibration."),
    ("0x00A4", "LORA_ADPT_PTR", "0x3F001200", "32-bit pointer targeting localized weight offset matrix."),
    ("0x00B0", "THRES_CONF_VAL", "0x0000FF7F", "Confidence score cut-off for object detection bounding boxes."),
    ("0x00C2", "ERR_LOG_CTRL", "0x00000001", "Trigger register for hardware diagnostic dump sequence.")
]
for addr, name, default, desc in registers:
    pdf.cell(30, 6, addr, 1, 0, 'C')
    pdf.cell(40, 6, name, 1, 0, 'L')
    pdf.cell(30, 6, default, 1, 0, 'C')
    pdf.cell(90, 6, desc, 1, 1, 'L')

pdf.output("data/raw/specs.pdf")


# ==========================================
# 2. 產生結構豐富的技術簡報 (PPTX)
# ==========================================
prs = Presentation()

# Slide 1: 封面
blank_layout = prs.slide_layouts[6] # 使用空白版面方便自訂元件
slide1 = prs.slides.add_slide(blank_layout)

title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(2.0))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Raydium Display AI Integration"
p.font.size = Pt(40)
p.font.bold = True

sub_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.5), Inches(1.5))
p2 = sub_box.text_frame.paragraphs[0]
p2.text = "Optimizing driver performance with MCP-based LLM workflows.\nAuthor: AI Engineering Team"
p2.font.size = Pt(18)

# Slide 2: 系統架構亮點 (條列式清單測試)
slide2 = prs.slides.add_slide(blank_layout)
header_box = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(0.8))
header_box.text_frame.paragraphs[0].text = "Key Architectural Highlights"
header_box.text_frame.paragraphs[0].font.size = Pt(28)
header_box.text_frame.paragraphs[0].font.bold = True

content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.0))
tf2 = content_box.text_frame
tf2.word_wrap = True

bullets = [
    "Edge-AI Upscaling: Integrated neural network pipes running locally on DDIC hardware layer.",
    "Model Context Protocol (MCP) Bindings: Exposing hardware register states to Claude Skills directly.",
    "FastAPI Gateway: High-concurrency RESTful interface allowing software agents to poll telemetry data.",
    "Dynamic Adapters: Support hot-swapping LoRA weights for specialized image optimization domains."
]
for i, bullet_text in enumerate(bullets):
    p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
    p.text = "• " + bullet_text
    p.font.size = Pt(14)
    p.space_after = Pt(14)

# Slide 3: 效能指標表格 (PPTX 表格測試)
slide3 = prs.slides.add_slide(blank_layout)
header_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(0.8))
header_box3.text_frame.paragraphs[0].text = "Inference Pipeline Performance Metrics"
header_box3.text_frame.paragraphs[0].font.size = Pt(28)
header_box3.text_frame.paragraphs[0].font.bold = True

# 建立表格：5 列 3 欄
table_shape = slide3.shapes.add_table(5, 3, Inches(0.8), Inches(1.8), Inches(8.4), Inches(3.5))
table = table_shape.table

# 設定欄寬
table.columns[0].width = Inches(3.4)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.5)

# 表格資料
table_data = [
    ["Resolution Target", "Latency (ms)", "Power Draw (mW)"],
    ["FHD+ (2400x1080) @ 60fps", "2.4 ms", "45 mW"],
    ["WQHD (3200x1440) @ 90fps", "4.1 ms", "78 mW"],
    ["UHD 4K (3840x2160) @ 120fps", "8.9 ms", "142 mW"],
    ["Super-Resolution Refinement Layer", "1.1 ms", "22 mW"]
]

for row_idx, row in enumerate(table_data):
    for col_idx, text in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        cell.text = text
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        if row_idx == 0:
            cell.text_frame.paragraphs[0].font.bold = True

prs.save("data/raw/presentation.pptx")

print("⚡ 內容更豐富的瑞鼎科技模擬資料（包含複雜表格、多章節、硬體描述）已產生於 data/raw/ !")