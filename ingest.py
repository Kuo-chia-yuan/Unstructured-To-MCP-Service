import os
import json
from pypdf import PdfReader
from pptx import Presentation

# ==========================================
# 📌 新增：切塊核心函式 (滑動視窗)
# ==========================================
def chunk_text(text: str, chunk_size: int = 100, chunk_overlap: int = 20) -> list:
    """
    將超長文本切成長度固定的區塊，並保留重疊範圍以維持語意連續性
    """
    chunks = []
    if not text:
        return chunks
        
    # 移除多餘的換行與空格，讓切塊更均勻
    text = " ".join(text.split())
    
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        # 每次前進「區塊大小 - 重疊大小」
        start += (chunk_size - chunk_overlap)
        
    return chunks

def process_files():
    raw_dir = "data/raw"
    processed_data = []

    if not os.path.exists(raw_dir):
        print(f"❌ 找不到目錄: {raw_dir}，請先執行 python make_data.py")
        return

    print("🚀 正在以 [精細切塊模式] 處理瑞鼎技術文件...")
    
    for filename in os.listdir(raw_dir):
        file_path = os.path.join(raw_dir, filename)
        content = ""
        
        try:
            if filename.endswith(".pdf"):
                print(f"正在讀取 PDF: {filename}")
                reader = PdfReader(file_path)
                content = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
            
            elif filename.endswith(".pptx"):
                print(f"正在讀取 PPTX: {filename}")
                prs = Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                content = " ".join(text_runs)
                
            # ==========================================
            # 📌 關鍵改動：在這裡將整份大文本切成多個小 Chunks
            # ==========================================
            if content:
                # 呼叫切塊函式：每個區塊 300 字，重疊 50 字
                text_chunks = chunk_text(content, chunk_size=100, chunk_overlap=20)
                
                # 每一段切出來的文字，各自存成知識庫的一筆獨立紀錄
                for idx, chunk in enumerate(text_chunks):
                    processed_data.append({
                        "source": filename,
                        "chunk_id": f"{filename}_chunk_{idx}", # 紀錄它是第幾塊，方便追蹤
                        "content": chunk,
                        "category": "IC_Spec" if "spec" in filename.lower() else "Technical_Slide"
                    })
                    
        except Exception as e:
            print(f"❌ 處理 {filename} 時出錯: {e}")

    # 儲存結果
    os.makedirs("data/processed", exist_ok=True)
    output_path = "data/processed/knowledge_base.json"
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(processed_data, f, ensure_ascii=False, indent=4)
    
    print(f"✅ 處理與切塊完成！共生成 {len(processed_data)} 個知識區塊。")
    print(f"結果已存至: {output_path}")

if __name__ == "__main__":
    process_files()